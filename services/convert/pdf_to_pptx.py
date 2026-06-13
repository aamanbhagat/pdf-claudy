"""PDF -> PowerPoint: one slide per page, each page rendered as a full-bleed image.

PDFs aren't slide decks, so the faithful, standard conversion is a picture per
slide (rather than guessing editable shapes). Uses PyMuPDF to rasterize.
"""
import os
import sys
import tempfile

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Emu

EMU_PER_PT = 12700  # 1pt = 12700 EMU

inp, out = sys.argv[1], sys.argv[2]
doc = fitz.open(inp)
prs = Presentation()
blank = prs.slide_layouts[6]

# One deck-wide slide size, taken from the first page (most PDFs are uniform).
first = doc[0].rect
prs.slide_width = Emu(int(first.width * EMU_PER_PT))
prs.slide_height = Emu(int(first.height * EMU_PER_PT))

zoom = 150 / 72  # 150 DPI
for page in doc:
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
    tmp = tempfile.mktemp(suffix=".png")
    pix.save(tmp)
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(tmp, 0, 0, width=prs.slide_width, height=prs.slide_height)
    os.remove(tmp)

prs.save(out)
